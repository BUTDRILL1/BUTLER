from __future__ import annotations

import time
from pathlib import Path
from typing import Any

from pydantic import BaseModel, Field

from butler.tools.base import Tool, ToolContext, ToolError


class SyncArgs(BaseModel):
    max_files: int = Field(default=2000, ge=1, le=20000)
    paths: list[str] | None = Field(default=None, description="Optional specific paths to sync incrementally.")


TEXT_EXTS = {".md", ".txt", ".csv"}
PDF_EXTS = {".pdf"}
OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}
ALL_INDEXABLE = TEXT_EXTS | PDF_EXTS | OFFICE_EXTS


def _read_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover - dependency missing is handled at runtime
        raise ToolError("PDF support requires the 'pypdf' package.") from e

    reader = PdfReader(str(path))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    return "\n".join(pages).strip()


def _read_office_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
        wb.close()
        return "\n".join(rows)
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
        return "\n".join(texts)
    return ""


def _extract_index_text(path: Path) -> tuple[str | None, bool]:
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTS:
        return path.read_text(encoding="utf-8", errors="replace"), False
    if suffix in PDF_EXTS:
        return _read_pdf_text(path), True
    if suffix in OFFICE_EXTS:
        try:
            return _read_office_text(path), False
        except Exception:
            return path.name, False  # At least index the filename
    return None, False


def _index_sync(ctx: ToolContext, args: SyncArgs) -> dict[str, Any]:
    if not ctx.sandbox.allowed_roots:
        raise ToolError("No allowed roots configured. Add one with /roots add <path>.")

    started = time.time()
    scanned = 0
    indexed = 0
    skipped_large = 0
    skipped_binary = 0

    if args.paths:
        # Incremental Sync
        targets = [Path(p) for p in args.paths]
    else:
        # Full Sync
        ctx.db.conn.execute("DELETE FROM files_fts;")
        # For memory, we don't wipe everything as it might contain chat history, 
        # but for simplicity in v1 full sync we just scan all.
        targets = []
        for root in ctx.sandbox.allowed_roots:
            targets.extend(list(root.rglob("*")))

    for path in targets:
        if indexed >= args.max_files:
            break
        if not path.is_file():
            continue
        scanned += 1
        if path.suffix.lower() not in ALL_INDEXABLE:
            skipped_binary += 1
            continue
        size = path.stat().st_size
        if size > ctx.config.max_file_bytes:
            skipped_large += 1
            continue
            
        content, is_pdf = _extract_index_text(path)
        if content is None:
            skipped_binary += 1
            continue
            
        if is_pdf and not content.strip():
            content = f"{path.name}"
            
        # Clear existing entries for this path to avoid duplicates
        ctx.db.conn.execute("DELETE FROM files_fts WHERE path = ?", (str(path),))
        # Note: Ideally we would also clear memory by metadata, but memory is additive in this v1.
        
        ctx.db.conn.execute(
            "INSERT INTO files_fts (path, content) VALUES (?, ?)",
            (str(path), content),
        )
        
        # if ctx.memory:
        #     chunks = [content[i:i+1000] for i in range(0, len(content), 800)]
        #     for i, chunk in enumerate(chunks):
        #         ctx.memory.add(
        #             chunk, 
        #             metadata={"path": str(path), "type": "file", "chunk_idx": i}
        #         )
        
        indexed += 1

    ctx.db.conn.commit()
    elapsed_ms = int((time.time() - started) * 1000)
    return {
        "mode": "incremental" if args.paths else "full",
        "scanned": scanned,
        "indexed": indexed,
        "elapsed_ms": elapsed_ms,
    }


def build() -> list[Tool]:
    return [
        Tool(
            name="index.sync",
            description="Index .md/.txt files under allowed roots into SQLite FTS.",
            input_model=SyncArgs,
            handler=_index_sync,
            side_effect=True,
        )
    ]
