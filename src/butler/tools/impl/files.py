from __future__ import annotations

import sqlite3
from pathlib import Path
from typing import Any

from pydantic import BaseModel, Field

from butler.sandbox import SandboxError
from butler.tools.base import Tool, ToolContext, ToolError


class ListArgs(BaseModel):
    path: str = Field(default=".")
    max_entries: int = Field(default=200, ge=1, le=5000)


def _list_dir(ctx: ToolContext, args: ListArgs) -> dict[str, Any]:
    if not ctx.sandbox.allowed_roots:
        raise ToolError("No allowed roots configured. Add one with /roots add <path>.")

    requested = Path(args.path)
    resolved = ctx.sandbox.ensure_allowed(requested if requested.is_absolute() else (ctx.sandbox.allowed_roots[0] / requested))
    if not resolved.exists():
        raise ToolError(f"Not found: {resolved}")
    if not resolved.is_dir():
        raise ToolError(f"Not a directory: {resolved}")

    entries: list[dict[str, Any]] = []
    for i, child in enumerate(sorted(resolved.iterdir(), key=lambda p: p.name.lower())):
        if i >= args.max_entries:
            break
        try:
            child_res = ctx.sandbox.ensure_allowed(child)
        except SandboxError:
            continue
        entries.append(
            {
                "name": child_res.name,
                "path": str(child_res),
                "is_dir": child_res.is_dir(),
                "size_bytes": child_res.stat().st_size if child_res.is_file() else None,
            }
        )
    return {"path": str(resolved), "entries": entries}


class ReadTextArgs(BaseModel):
    path: str


def _read_pdf(p: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(p))
    pages: list[str] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i+1} ---\n{text.strip()}")
    return "\n\n".join(pages) if pages else "(No extractable text found in this PDF)"


def _read_docx(p: Path) -> str:
    from docx import Document
    doc = Document(str(p))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(paragraphs) if paragraphs else "(No text found in this Word document)"


def _read_xlsx(p: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(p), read_only=True, data_only=True)
    sheets: list[str] = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"--- Sheet: {name} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets) if sheets else "(No data found in this Excel file)"


def _read_pptx(p: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(p))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "(No text found in this PowerPoint)"


def _read_csv(p: Path) -> str:
    import csv
    with open(p, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = [" | ".join(row) for row in reader]
    return "\n".join(rows[:500]) if rows else "(Empty CSV file)"


_FORMAT_READERS: dict[str, Any] = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
    ".csv": _read_csv,
}


def _read_text(ctx: ToolContext, args: ReadTextArgs) -> dict[str, Any]:
    if not ctx.sandbox.allowed_roots:
        raise ToolError("No allowed roots configured. Add one with /roots add <path>.")
    p = ctx.sandbox.ensure_allowed(Path(args.path))
    if not p.exists() or not p.is_file():
        raise ToolError(f"Not a file: {p}")
    size = p.stat().st_size
    if size > ctx.config.max_file_bytes:
        raise ToolError(f"File too large ({size} bytes). Limit is {ctx.config.max_file_bytes}.")

    ext = p.suffix.lower()
    reader_fn = _FORMAT_READERS.get(ext)
    if reader_fn:
        try:
            content = reader_fn(p)
            return {"path": str(p), "size_bytes": size, "format": ext, "content": content}
        except Exception as e:
            return {"path": str(p), "size_bytes": size, "format": ext, "error": f"Failed to read {ext} file: {e}"}

    content = p.read_text(encoding="utf-8", errors="replace")
    return {"path": str(p), "size_bytes": size, "content": content}


class SearchArgs(BaseModel):
    query: str = Field(min_length=1, max_length=400)
    # Allow higher inputs but clamp internally (Gemma-friendly).
    limit: int = Field(default=10, ge=1, le=10_000)


def _clean_snippet(s: str) -> str:
    return s.strip().replace("\n", " ")


def _truncate(s: str, max_chars: int) -> str:
    if len(s) <= max_chars:
        return s
    return s[:max_chars].rstrip()


def _is_invalid_fts_query_error(e: sqlite3.OperationalError) -> bool:
    msg = str(e).lower()
    return ("fts" in msg and "syntax" in msg) or ("syntax error" in msg) or ("parse" in msg)


def _is_bm25_missing_error(e: sqlite3.OperationalError) -> bool:
    msg = str(e).lower()
    return "bm25" in msg and ("no such function" in msg or "unknown function" in msg or "bm25" in msg)


def _files_search(ctx: ToolContext, args: SearchArgs) -> dict[str, Any]:
    limit = max(1, min(int(args.limit), 10))
    query = args.query.strip()
    
    # 1. Path / Filename Search (most common use-case)
    path_results = []
    try:
        rows = ctx.db.conn.execute(
            "SELECT path FROM files_fts WHERE path LIKE ? LIMIT ?",
            (f"%{query}%", limit),
        ).fetchall()
        for r in rows:
            path_results.append({
                "path": r["path"],
                "score": 2.0,
                "snippet": f"Filename match: {r['path']}",
                "method": "filename"
            })
    except sqlite3.OperationalError:
        pass

    # 2. Semantic Search (Phase 5 Upgrade)
    semantic_results = []
    try:
        if ctx.memory:
            mem_hits = ctx.memory.search(query, limit=limit)
            for hit in mem_hits:
                if hit["metadata"].get("type") == "file":
                    semantic_results.append({
                        "path": hit["metadata"]["path"],
                        "score": hit["score"],
                        "snippet": _clean_snippet(_truncate(hit["content"], 200)),
                        "method": "semantic"
                    })
    except Exception as e:
        print(f"Semantic Search Warning: {e}")

    # 3. Traditional FTS Content Search
    fts_results = []
    try:
        rows = ctx.db.conn.execute(
            """
            SELECT path, snippet(files_fts, 1, '', '', ' … ', 8) AS snippet
            FROM files_fts WHERE files_fts MATCH ? LIMIT ?
            """,
            (query, limit),
        ).fetchall()
        for r in rows:
            fts_results.append({
                "path": r["path"],
                "score": 1.0,
                "snippet": _clean_snippet(_truncate(r["snippet"] or "", 200)),
                "method": "keyword"
            })
    except sqlite3.OperationalError:
        pass

    # 4. Merge and Deduplicate (path matches first)
    all_results = path_results + semantic_results + fts_results
    seen_paths = set()
    final_results = []
    for res in all_results:
        if res["path"] not in seen_paths:
            final_results.append(res)
            seen_paths.add(res["path"])
            
    return {
        "query": query, 
        "results": final_results[:limit]
    }


def build() -> list[Tool]:
    return [
        Tool(
            name="files.list",
            description="List directory entries within allowed roots.",
            input_model=ListArgs,
            handler=_list_dir,
        ),
        Tool(
            name="files.read_text",
            description="Read a UTF-8 text file within allowed roots (size-capped).",
            input_model=ReadTextArgs,
            handler=_read_text,
        ),
        Tool(
            name="files.search",
            description="Search indexed .md/.txt content (requires /index sync). Returns <=10 short snippets.",
            input_model=SearchArgs,
            handler=_files_search,
            side_effect=False,
        ),
    ]
